import os
import re
import logging
import pandas as pd
import PyPDF2
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path

logger = logging.getLogger(__name__)


# -----------------------------
# FILE TYPE HANDLER
# -----------------------------
def extract_text_from_file(file_path: str) -> str:
    """Auto-detect file type and extract text with proper cleaning."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        text = extract_pdf_text(file_path)
        if not text.strip():  # fallback to OCR
            text = extract_text_from_pdf_ocr(file_path)
        return text

    elif ext in [".ppt", ".pptx"]:
        return extract_ppt_text(file_path)

    elif ext in [".xls", ".xlsx"]:
        return extract_excel_text(file_path)

    else:
        logger.warning(f"Unsupported file type: {ext}")
        return ""


# -----------------------------
# LIGHT CLEANING (SAFE)
# -----------------------------
def clean_text(text: str) -> str:
    """
    Minimal cleaning:
    - Remove bullets
    - Normalize whitespace
    - Keep punctuation, uppercase, structure
    """
    text = text.replace("\u2022", "")  # bullet character
    text = re.sub(r"[ ]{2,}", " ", text)  # collapse multiple spaces
    text = re.sub(r"\n{3,}", "\n\n", text)  # limit blank lines
    return text.strip()


# -----------------------------
# PDF EXTRACTION
# -----------------------------
def extract_pdf_text(file_path: str) -> str:
    """Extract text from normal PDF."""
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            for page in reader.pages:
                content = page.extract_text() or ""
                pages.append(content)

        full_text = "\n\n--- PAGE BREAK ---\n\n".join(pages)
        return clean_text(full_text)

    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return ""


def extract_text_from_pdf_ocr(pdf_path: str) -> str:
    """Extract text from scanned PDF using OCR."""
    try:
        logger.info(f"Starting OCR for scanned PDF: {pdf_path}")
        images = convert_from_path(pdf_path)

        text_chunks = []
        for idx, img in enumerate(images):
            logger.info(f"OCR processing page {idx + 1}/{len(images)}")
            page_text = pytesseract.image_to_string(img)
            text_chunks.append(clean_text(page_text))

        full_text = "\n\n--- PAGE OCR BREAK ---\n\n".join(text_chunks)
        return full_text.strip()

    except Exception as e:
        logger.error(f"OCR extraction error: {e}")
        return ""


# -----------------------------
# PPT EXTRACTION (Improved)
# -----------------------------
def extract_ppt_text(file_path: str) -> str:
    """
    Extract text from PPTX while preserving:
    - slide titles
    - paragraphs
    - bullet points
    - slide boundaries
    """
    try:
        prs = Presentation(file_path)
        slides_output = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_lines = []

            # Slide title
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
                slide_lines.append(f"### Slide {slide_index}: {title}")

            # Other slide text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if shape != slide.shapes.title:
                        slide_lines.append(text)

            slide_block = "\n".join(slide_lines)
            slides_output.append(slide_block)

        full_text = "\n\n--- SLIDE BREAK ---\n\n".join(slides_output)
        return clean_text(full_text)

    except Exception as e:
        logger.error(f"PPT extraction failed: {e}")
        return ""


# -----------------------------
# EXCEL EXTRACTION
# -----------------------------
def extract_excel_text(file_path: str) -> str:
    """Extract text from Excel by flattening rows."""
    try:
        df = pd.read_excel(file_path)
        df = df.dropna(how="all").fillna("")
        text = " ".join(df.astype(str).values.flatten())
        return clean_text(text)

    except Exception as e:
        logger.error(f"Excel extraction failed: {e}")
        return ""


# -----------------------------
# CHUNKING (RAG Optimized)
# -----------------------------
def chunk_text(text: str, chunk_size: int = 200) -> list[str]:
    """
    Chunk text into ~200-word blocks for best embedding performance.
    """
    words = text.split()
    chunks = [
        " ".join(words[i:i + chunk_size])
        for i in range(0, len(words), chunk_size)
    ]
    return chunks
